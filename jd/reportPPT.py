from pptx.util import Pt
from pptx.util import Inches


def text_ppt(prs, title):
    # 插入一页幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    # 对ppt的修改
    body_shape = slide.shapes.placeholders  # body_shape为本页ppt中所有shapes
    # 在文本框中添加新段落
    # 在创建的这页幻灯片文本框中添加文字
    title_shape = slide.shapes.title  # 取本页ppt的title
    title_shape.text = title  # 向title文本框写如文字
    return body_shape


def add_para(prs, title, **paragraphs):
    body_shape = text_ppt(prs, title)
    i = 0
    for paragraph in paragraphs:
        i += 1
        new_paragraph = body_shape[1].text_frame.add_paragraph()  # 在第二个shape中的文本框架中添加新段落
        new_paragraph.text = paragraphs[paragraph]  # 新段落中文字
        new_paragraph.font.italic = False  # 文字斜体
        if i == 1:
            new_paragraph.font.bold = False  # 文字加粗
            new_paragraph.font.size = Pt(20)  # 文字大小
            new_paragraph.font.underline = False  # 文字下划线
        else:
            new_paragraph.font.bold = False  # 文字加粗
            new_paragraph.font.size = Pt(16)
            new_paragraph.font.underline = True  # 文字下划线
        new_paragraph.level = 1  # 新段落的级别


def add_pic(prs, title, picpath):
    # 插入一页幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # 对ppt的修改
    title_shape = slide.shapes.title  # 取本页ppt的title
    title_shape.text = title  # 向title文本框写如文字
    # 添加图片
    img_path = picpath  # 文件路径
    left, top, width, height = Inches(1), Inches(1.5), Inches(8), Inches(6)  # 预设位置及大小
    pic = slide.shapes.add_picture(img_path, left, top, width, height)  # 在指定位置按预设值添加图片
